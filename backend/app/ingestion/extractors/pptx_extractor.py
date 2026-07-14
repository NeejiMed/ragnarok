from pptx import Presentation


def extract_pptx(file_path: str) -> list[dict]:
    """
    Extracts text from a PPTX file, one Document per slide.
    Includes all shape text frames + speaker notes.
    """
    prs = Presentation(file_path)
    results = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = _extract_slide_text(slide)
        notes_text = _extract_notes(slide)

        combined = slide_text
        if notes_text:
            combined += f"\n\n[Speaker notes]\n{notes_text}"

        results.append(
            {"content": combined.strip(), "page": slide_number, "extraction_method": "native"}
        )

    return results


def _extract_slide_text(slide) -> str:
    """
    Extracts text from all shapes in a slide.
    """
    texts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            line = "".join(run.text for run in paragraph.runs).strip()
            if line:
                texts.append(line)
    return "\n".join(texts)


def _extract_notes(slide) -> str:
    """
    Extracts speaker notes from a slide.
    """
    if not slide.has_notes_slide:
        return ""
    notes_slide = slide.notes_slide
    notes_frame = notes_slide.notes_text_frame
    if not notes_frame:
        return ""
    return "\n".join(
        paragraph.text.strip() for paragraph in notes_frame.paragraphs if paragraph.text.strip()
    )
