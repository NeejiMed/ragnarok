from pathlib import Path
from pptx import Presentation
from pptx.exc import PackageNotFoundError
from pptx.util import Inches
import pytest

from backend.app.ingestion.extractors.pptx_extractor import extract_pptx

def make_text_pptx(path: str, slides: list[str]) -> None:
    """One slide per string, text in a title box."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout
    for text in slides:
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(5)
        )
        txBox.text_frame.text = text
    prs.save(path)

def make_pptx_with_notes(path: str, slide_text: str, notes_text: str) -> None:
    """Single slide with content and speaker notes."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(5)
    )
    txBox.text_frame.text = slide_text
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text # type: ignore
    prs.save(path)

def test_single_slide_pptx_extraction():
    pptx_path = Path(__file__).parent / "fixtures" / "single-slide.pptx"
    #make_text_pptx(str(pptx_path), ["Hello, World!"])
    extracted_content = extract_pptx(str(pptx_path))
    print(extracted_content)

    assert len(extracted_content) == 1
    assert extracted_content[0]["extraction_method"] == "native"
    assert extracted_content[0]["page"] is 1
    assert "PowerPointA" in extracted_content[0]["content"]

def test_pptxx_multiple_slides_extraction():
    pptx_path = Path(__file__).parent / "fixtures" / "multi-slide.pptx"
    #make_text_pptx(str(pptx_path), ["Slide 1", "Slide 2", "Slide 3"])
    extracted_content = extract_pptx(str(pptx_path))
    print(extracted_content)

    assert len(extracted_content) == 3
    for i, item in enumerate(extracted_content, start=1):
        assert item["page"] == i
        assert item["extraction_method"] == "native"

def test_pptx_with_notes_extraction():
    pptx_path = Path(__file__).parent / "fixtures" / "slide-with-notes.pptx"
    #make_pptx_with_notes(str(pptx_path), "Slide Content", "Speaker Notes")
    extracted_content = extract_pptx(str(pptx_path))
    print(extracted_content)

    assert len(extracted_content) == 3
    assert extracted_content[0]["extraction_method"] == "native"
    assert extracted_content[0]["page"] == 1
    assert "[Speaker notes]" in extracted_content[0]["content"]

def test_invalid_path_pptx_extraction():
    fake_path = "backend/tests/does_not_exist.pptx"
    with pytest.raises(PackageNotFoundError):
        extract_pptx(fake_path)